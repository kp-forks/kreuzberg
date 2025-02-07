from __future__ import annotations

import re
from asyncio import gather
from contextlib import suppress
from html import escape
from io import BytesIO
from typing import TYPE_CHECKING, Any, cast

import tesserocr
from anyio import Path as AsyncPath
from charset_normalizer import detect
from typing_extensions import Never

from kreuzberg._mime_types import PANDOC_MIME_TYPE_EXT_MAP
from kreuzberg._string import normalize_spaces, safe_decode
from kreuzberg._sync import run_sync
from kreuzberg.exceptions import MissingDependencyError, ParsingError

try:
    import pypdfium2
except ImportError:  # pragma: no cover
    pypdfium2 = Never

try:
    import pypandoc
except ImportError:  # pragma: no cover
    pypandoc = Never

try:
    import pptx
    import pptx.enum.shapes
except ImportError:  # pragma: no cover
    pptx = Never  # type: ignore[assignment]
    shapes = Never

try:
    import html_to_markdown
except ImportError:  # pragma: no cover
    html_to_markdown = Never  # type: ignore[assignment]

if TYPE_CHECKING:  # pragma: no cover
    from pathlib import Path

    from PIL.Image import Image


def _assert_dependency_installed(name: str, dependency: Any) -> None:
    """Assert that a dependency is installed.

    Args:
        name: The name of the dependency.
        dependency: The dependency to check.

    Raises:
        MissingDependencyError: If the dependency is not installed

    Returns:
        None
    """
    if dependency is Never:
        raise MissingDependencyError(f"missing dependency {name}")


def _convert_pdf_to_images(file_path: Path) -> list[Image]:
    """Convert a PDF file to images.

    Args:
        file_path: The path to the PDF file.

    Returns:
        A list of paths to the images.
    """
    _assert_dependency_installed("pdfium2", pypdfium2)

    try:
        pdf = pypdfium2.PdfDocument(str(file_path))
        return [page.render(scale=2.0).to_pil() for page in pdf]
    except pypdfium2.PdfiumError as e:
        raise ParsingError(
            "Could not convert PDF to images", context={"file_path": str(file_path), "error": str(e)}
        ) from e


async def _ocr_image_with_tesseract(image: Image) -> str:
    """Perform OCR on an image using tesserocr.

    Args:
        image: The image to perform OCR on.

    Returns:
        The extracted text.
    """
    _assert_dependency_installed("tesserocr", tesserocr)

    try:
        result = await run_sync(tesserocr.image_to_text, image, psm=tesserocr.PSM.AUTO_OSD, lang="eng")
    except RuntimeError as e:
        raise ParsingError("Could not extract text from image", context={"error": str(e)}) from e

    return normalize_spaces(result)


async def _extract_pdf_with_tesseract(file_path: Path) -> str:
    """Extract text from a scanned PDF file using pytesseract.

    Args:
        file_path: The path to the PDF file.

    Raises:
        ParsingError: If the text could not be extracted from the PDF file.

    Returns:
        The extracted text.
    """
    _assert_dependency_installed("tesserocr", tesserocr)

    ocr_results = await gather(*[_ocr_image_with_tesseract(img) for img in _convert_pdf_to_images(file_path)])
    return normalize_spaces("\n".join(ocr_results))


def _extract_pdf_with_pdfium2(file_path: Path) -> str:
    """Extract text from a searchable PDF file using pypdfium2.

    Args:
        file_path: The path to the PDF file.

    Raises:
        ParsingError: If the text could not be extracted from the PDF file.

    Returns:
        The extracted text.
    """
    _assert_dependency_installed("pdfium2", pypdfium2)
    try:
        document = pypdfium2.PdfDocument(file_path)
        text = "\n".join(page.get_textpage().get_text_range() for page in document)
        return normalize_spaces(text)
    except pypdfium2.PdfiumError as e:
        # TODO: add test case
        raise ParsingError(
            "Could not extract text from PDF file", context={"file_path": str(file_path), "error": str(e)}
        ) from e


async def _extract_pdf_file(file_path: Path, force_ocr: bool = False) -> str:
    """Extract text from a PDF file.

    Args:
        file_path: The path to the PDF file.
        force_ocr: Whether or not to force OCR on PDF files that have a text layer. Default = false.

    Returns:
        The extracted text.
    """
    if not force_ocr and (content := await run_sync(_extract_pdf_with_pdfium2, file_path)):
        return normalize_spaces(content)

    return await _extract_pdf_with_tesseract(file_path)


async def _extract_content_with_pandoc(file_data: bytes, mime_type: str, encoding: str | None = None) -> str:
    """Extract text using pandoc.

    Args:
        file_data: The content of the file.
        mime_type: The mime type of the file.
        encoding: An optional encoding to use when decoding the string.

    Raises:
        ParsingError: If the text could not be extracted from the file using pandoc.

    Returns:
        The extracted text.
    """
    _assert_dependency_installed("pypandoc", pypandoc)

    ext = PANDOC_MIME_TYPE_EXT_MAP[mime_type]
    encoding = encoding or detect(file_data)["encoding"] or "utf-8"
    try:
        return normalize_spaces(
            cast(str, await run_sync(pypandoc.convert_text, file_data, to="md", format=ext, encoding=encoding))
        )
    except RuntimeError as e:
        # TODO: add test case
        raise ParsingError(
            f"Could not extract text from {PANDOC_MIME_TYPE_EXT_MAP[mime_type]} file contents",
            context={"error": str(e)},
        ) from e


async def _extract_file_with_pandoc(file_path: Path | str, mime_type: str) -> str:
    """Extract text using pandoc.

    Args:
        file_path: The path to the file.
        mime_type: The mime type of the file.

    Raises:
        ParsingError: If the text could not be extracted from the file using pandoc.

    Returns:
        The extracted text.
    """
    _assert_dependency_installed("pandoc", pypandoc)

    ext = PANDOC_MIME_TYPE_EXT_MAP[mime_type]
    try:
        return normalize_spaces(cast(str, await run_sync(pypandoc.convert_file, file_path, to="md", format=ext)))
    except RuntimeError as e:
        raise ParsingError(
            f"Could not extract text from {PANDOC_MIME_TYPE_EXT_MAP[mime_type]} file",
            context={"file_path": str(file_path), "error": str(e)},
        ) from e


async def _extract_image_with_tesseract(file_path: Path | str) -> str:
    """Extract text from an image file.

    Args:
        file_path: The path to the image file.

    Raises:
        ParsingError: If the text could not be extracted from the image file.

    Returns:
        The extracted content.
    """
    _assert_dependency_installed("tesserocr", tesserocr)

    try:
        result = normalize_spaces(cast(str, tesserocr.file_to_text(str(file_path))))
        return result
    except RuntimeError as e:
        raise ParsingError(
            "Could not extract text from image file", context={"file_path": str(file_path), "error": str(e)}
        ) from e


async def _extract_pptx_file(file_path_or_contents: Path | bytes) -> str:
    """Extract text from a PPTX file.

    Notes:
        This function is based on code vendored from `markitdown`, which has an MIT license as well.

    Args:
        file_path_or_contents: The path to the PPTX file or its contents as bytes.

    Returns:
        The extracted text content
    """
    _assert_dependency_installed("pptx", pptx)

    md_content = ""
    file_contents = (
        file_path_or_contents
        if isinstance(file_path_or_contents, bytes)
        else await AsyncPath(file_path_or_contents).read_bytes()
    )
    presentation = pptx.Presentation(BytesIO(file_contents))

    for index, slide in enumerate(presentation.slides):
        md_content += f"\n\n<!-- Slide number: {index + 1} -->\n"

        title = slide.shapes.title

        for shape in slide.shapes:
            if shape.shape_type == shapes.MSO_SHAPE_TYPE.PICTURE or (
                shape.shape_type == shapes.MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image")
            ):
                alt_text = ""
                with suppress(AttributeError):
                    # access non-visual properties
                    alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")  # noqa: SLF001

                filename = re.sub(r"\W", "", shape.name) + ".jpg"
                md_content += f"\n![{alt_text if alt_text else shape.name}]({filename})\n"

            elif shape.shape_type == shapes.MSO_SHAPE_TYPE.TABLE:
                html_table = "<table>"
                first_row = True

                for row in shape.table.rows:
                    html_table += "<tr>"

                    for cell in row.cells:
                        tag = "th" if first_row else "td"
                        html_table += f"<{tag}>{escape(cell.text)}</{tag}>"

                    html_table += "</tr>"
                    first_row = False

                html_table += "</table>"
                md_content += "\n" + html_table + "\n"

            elif shape.has_text_frame:
                md_content += "# " + shape.text.lstrip() + "\n" if shape == title else shape.text + "\n"

        md_content = md_content.strip()
        if slide.has_notes_slide:
            md_content += "\n\n### Notes:\n"
            notes_frame = slide.notes_slide.notes_text_frame

            if notes_frame is not None:
                md_content += notes_frame.text

            md_content = md_content.strip()

    return normalize_spaces(md_content)


async def _extract_html_string(file_path_or_contents: Path | bytes) -> str:
    """Extract text from an HTML string.

    Args:
        file_path_or_contents: The HTML content.

    Returns:
        The extracted text content.
    """
    _assert_dependency_installed("html_to_markdown", html_to_markdown)
    content = (
        safe_decode(file_path_or_contents)
        if isinstance(file_path_or_contents, bytes)
        else await AsyncPath(file_path_or_contents).read_text()
    )
    return normalize_spaces(await run_sync(html_to_markdown.convert_to_markdown, content))
